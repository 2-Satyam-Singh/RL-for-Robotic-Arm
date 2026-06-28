from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.33)   # 16:9 aspect ratio
prs.slide_height = Inches(7.5)

# ==================== SLIDE 1: Overall Training Workflow ====================
slide1 = prs.slides.add_slide(prs.slide_layouts[5])
title1 = slide1.shapes.title
title1.text = "Overall Training Workflow (Horizontal View)"
title1.text_frame.paragraphs[0].font.size = Pt(28)

# Function to add nice block
def add_block(slide, text, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*color)
    shape.line.color.rgb = RGBColor(0, 0, 0)
    shape.line.width = Pt(2)
    tx = shape.text_frame
    tx.text = text
    tx.paragraphs[0].font.size = Pt(14)
    tx.paragraphs[0].font.bold = True
    tx.paragraphs[0].alignment = 1  # center
    return shape

# Blocks (Horizontal layout)
blocks = [
    ("Start Training\ntrain.py", 0.5, 1.2, 1.8, 1.2, (0, 102, 204)),
    ("PandaEnv.reset()", 2.8, 1.2, 1.8, 1.2, (0, 153, 76)),
    ("PandaController\nreset() + reward reset", 5.1, 1.2, 2.0, 1.2, (0, 102, 204)),
    ("Normalized Obs\n(joints + entity)", 7.6, 1.2, 1.8, 1.2, (255, 140, 0)),
    ("PPO.select_action()\nActor-Critic + tanh", 10.0, 1.2, 2.2, 1.2, (204, 0, 102)),
]

for i, (txt, l, t, w, h, col) in enumerate(blocks):
    add_block(slide1, txt, Inches(l), Inches(t), Inches(w), Inches(h), col)

# Arrows
arrow_positions = [(2.3, 1.8), (4.8, 1.8), (7.1, 1.8), (9.4, 1.8)]
for x, y in arrow_positions:
    arrow = slide1.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(0.4), Inches(0.3))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(0, 0, 0)

# Bottom row (loop)
add_block(slide1, "store_transition()\n+ PPO.learn()", Inches(3.0), Inches(4.0), Inches(2.5), Inches(1.0), (0, 153, 76))
add_block(slide1, "Sparse Reward\n(+100 touch / +1000 success)", Inches(6.0), Inches(4.0), Inches(2.8), Inches(1.0), (204, 0, 102))
add_block(slide1, "Logger\nCSV + Plot", Inches(9.5), Inches(4.0), Inches(2.0), Inches(1.0), (255, 140, 0))

# ==================== SLIDE 2: PPO Algorithm ====================
slide2 = prs.slides.add_slide(prs.slide_layouts[5])
title2 = slide2.shapes.title
title2.text = "PPO Algorithm Flow (Horizontal View)"
title2.text_frame.paragraphs[0].font.size = Pt(28)

blocks2 = [
    ("Start Episode", 0.8, 2.0, 1.8, 1.2, (0, 102, 204)),
    ("Collect Rollout\n2048 steps", 3.0, 2.0, 2.0, 1.2, (0, 153, 76)),
    ("Compute GAE\n+ Returns", 5.5, 2.0, 2.0, 1.2, (255, 140, 0)),
    ("For 8 epochs", 8.0, 2.0, 1.8, 1.2, (204, 0, 102)),
    ("Shuffle Minibatches\n(size 64)", 10.3, 2.0, 2.2, 1.2, (0, 102, 204)),
]

for txt, l, t, w, h, col in blocks2:
    add_block(slide2, txt, Inches(l), Inches(t), Inches(w), Inches(h), col)

# Final block
add_block(slide2, "Adam Update\nActor-Critic", Inches(4.0), Inches(4.5), Inches(2.5), Inches(1.0), (0, 153, 76))
add_block(slide2, "Next Episode\n(10,000 total)", Inches(7.5), Inches(4.5), Inches(2.5), Inches(1.0), (255, 140, 0))

print("✅ Slides created successfully! File saved as: RL_Panda_Paper_Slides.pptx")
prs.save("RL_Panda_Paper_Slides.pptx")