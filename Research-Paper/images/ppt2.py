from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs, title, subtitle, author):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    subtitle_shape.text = f"{subtitle}\n\n{author}"
    
    # Speaker Notes
    notes = slide.notes_slide.notes_text_frame
    notes.text = "Good morning, members of the committee. I am Satyam Singh. My research addresses the problem of morphological rigidity in robotic control. Specifically, we propose a framework for the automated synthesis of implicit inverse kinematic (IK) solvers across heterogeneous serial manipulators."

def add_content_slide(prs, title, bullets, notes_text, image_configs=None):
    """
    image_configs: list of dicts {'text': 'Placeholder Text', 'left': Inches(x), 'top': Inches(y), 'width': Inches(w), 'height': Inches(h)}
    """
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Blank slide with title
    
    # Set Title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    
    # Add Bullets
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(4.5) if image_configs else Inches(9.0)
    height = Inches(5.0)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = bullet
        p.level = 0
        p.font.size = Pt(22)
        p.space_after = Pt(14)
    
    # Add Image Placeholders (Grey Boxes with Text)
    if image_configs:
        for img in image_configs:
            shape = slide.shapes.add_shape(
                1, # Rectangle
                img['left'], img['top'], img['width'], img['height']
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(200, 200, 200)
            
            p = shape.text_frame.paragraphs[0]
            p.text = img['text']
            p.alignment = PP_ALIGN.CENTER
            p.font.color.rgb = RGBColor(0, 0, 0)
            
    # Add Speaker Notes
    notes = slide.notes_slide.notes_text_frame
    notes.text = notes_text

# Initialize Presentation
prs = Presentation()

# --- SLIDE 1: Title (Context) ---
add_title_slide(
    prs, 
    "Adaptive Determination of Robot Structure\nand Inverse Kinematics through RL", 
    "Implicit IK Learning in Contact-Rich Environments", 
    "Satyam Singh\nDepartment of Mechanical Engineering, MANIT Bhopal"
)

# --- SLIDE 2: Problem ---
add_content_slide(
    prs,
    "The Paradigm of Morphological Generalization",
    [
        "Modern robots perform dynamic tasks, but control laws are 'body-locked'.",
        "Modifying structural topology (link lengths, DoF) obsoletes the analytical model.",
        "Goal: Synthesize control policies dynamically to adapt to arbitrary physical constraints."
    ],
    "We see robots like this quadruped performing dynamic tasks, but their control laws are typically 'body-locked.' If link lengths change or joints are added, the math breaks. My work provides a 'Unified Factory' that builds a custom control brain for any robot body it is given.",
    [{'text': 'INSERT: Dog.png\n(Visual: Quadruped)', 'left': Inches(5.2), 'top': Inches(2.0), 'width': Inches(4.5), 'height': Inches(3.5)}]
)

# --- SLIDE 3: Problem ---
add_content_slide(
    prs,
    "Analytical IK and the Jacobian Bottleneck",
    [
        "Traditional IK relies on the Jacobian Pseudoinverse: $\Delta \theta = J^{\dagger} \Delta x$",
        "Analytical methods suffer numerical instability near kinematic singularities.",
        "We propose treating IK as an implicit function approximation problem, learned through interaction."
    ],
    "Standard IK relies on the inversion of the Jacobian matrix. For MSMEs building custom robots, this is a bottleneck. Analytical models are brittle near singularities and expensive to re-derive every time a link length is modified. We need an Implicit IK Solver.",
    [{'text': 'INSERT: Diagram of Singular Configuration', 'left': Inches(5.2), 'top': Inches(2.0), 'width': Inches(4.5), 'height': Inches(3.5)}]
)

# --- SLIDE 4: Solution ---
add_content_slide(
    prs,
    "Unified Pipeline for Controller Synthesis",
    [
        "Automated parsing of Simulation Description Format (SDF) files.",
        "Dynamic instantiation of Neural Network layers based on joint topology.",
        "A structure-agnostic architecture: identical codebase for 3-DoF to 7-DoF."
    ],
    "This is my core innovation: the Automated Pipeline. When a new robot description is loaded, the system identifies the joint count and dynamically instantiates the Neural Network layers. This isn't one brain for all robots; it's a factory that builds a specific brain for the robot you provide.",
    [{'text': 'INSERT: System Flowchart\nSDF -> Parser -> NN -> Gazebo', 'left': Inches(5.2), 'top': Inches(2.0), 'width': Inches(4.5), 'height': Inches(3.5)}]
)

# --- SLIDE 5: Solution ---
add_content_slide(
    prs,
    "Policy Optimization and Stochastic Control",
    [
        "Continuous control via Proximal Policy Optimization (PPO).",
        "Clipped Surrogate Objective: $L^{CLIP}(\phi) = \hat{\mathbb{E}}_t [\min(r_t(\phi) \hat{A}_t, \text{clip}(...) \hat{A}_t)]$",
        "Action Denormalization maps bounded outputs to physical joint limits: $a_{real} = l_{low} + \frac{(a_{norm} + 1)}{2} (l_{high} - l_{low})$"
    ],
    "We utilize Proximal Policy Optimization to ensure stable policy updates. By clipping the probability ratio, we avoid large, destabilizing updates that would violate mechanical limits. By normalizing all actions, the RL agent treats different motor torques through a standardized control lens.",
    [{'text': 'INSERT: PPO Actor-Critic Diagram', 'left': Inches(5.5), 'top': Inches(2.0), 'width': Inches(4.0), 'height': Inches(3.0)}]
)

# --- SLIDE 6: Solution ---
add_content_slide(
    prs,
    "The Sparse Reward Hypothesis",
    [
        "Dense distance shaping induces 'safe hovering' local optima to avoid contact physics.",
        "Proposed 'Massive Boost' Strategy: $R = 100(\text{Contact}) + 1000(\text{Success})$",
        "Forces the agent to endure initial exploration variance to prioritize high-impact force."
    ],
    "This was my most significant discovery. In contact physics, 'dense' rewards fail because the robot becomes 'scared' of the unpredictable forces of collision. I implemented a Massive Boost sparse reward. By providing a +1000 signal only on success, the agent learns to prioritize the high-impact force needed.",
    [{'text': 'INSERT: Reward Graph/Equation Breakdown', 'left': Inches(5.2), 'top': Inches(2.0), 'width': Inches(4.5), 'height': Inches(3.5)}]
)

# --- SLIDE 7: Evidence ---
add_content_slide(
    prs,
    "Cross-Morphology Implementation",
    [
        "Testing the 'Adaptive Determination' hypothesis across varied DoF.",
        "3-DoF Custom Arm: High rotational inertia on Link 2 ($0.097 kg \cdot m^2$).",
        "7-DoF Panda Arm: High redundancy ($18 kg$ total mass)."
    ],
    "To validate the pipeline's adaptability, I performed a comparative study across multiple morphologies. We aren't training in a vacuum; the environment accounts for Real-World Mechanical Constraints like static friction. The PPO agent must learn to compensate for varying momentum distributions.",
    [
        {'text': '3DOF.png', 'left': Inches(0.5), 'top': Inches(4.0), 'width': Inches(2.5), 'height': Inches(2.5)},
        {'text': '5DOF.png', 'left': Inches(3.5), 'top': Inches(4.0), 'width': Inches(2.5), 'height': Inches(2.5)},
        {'text': 'Panda.png', 'left': Inches(6.5), 'top': Inches(4.0), 'width': Inches(2.5), 'height': Inches(2.5)}
    ]
)

# --- SLIDE 8: Evidence ---
add_content_slide(
    prs,
    "Convergence Results: Sparse PPO (7-DoF)",
    [
        "Task: Non-prehensile pushing (off-table drop).",
        "Achieved a stable convergence profile over 10,000 episodes.",
        "Agent utilizes distal joints to maintain contact while proximal joints provide thrust ('sweeping')."
    ],
    "The 7-DoF Franka Panda achieved a 51% success rate within a 10,000-episode budget. The convergence profile indicates that after an initial high-variance exploration phase, the agent discovers the sparse reward manifold. It learned a sweeping motion, utilizing its redundancy to maintain object contact.",
    [{'text': 'INSERT: sparse_ppo_10k.png', 'left': Inches(5.2), 'top': Inches(2.0), 'width': Inches(4.5), 'height': Inches(3.5)}]
)

# --- SLIDE 9: Evidence ---
add_content_slide(
    prs,
    "Ablation Study: Algorithmic & Shaping Failures",
    [
        "Discrete DQN: Failed entirely. Curse of dimensionality ($3^7 = 2,187$ combinations/step).",
        "Dense PPO: Stagnated near zero. Confirms the 'safe hovering' hypothesis.",
        "Validates continuous action spaces and sparse signals for implicit IK."
    ],
    "I compared PPO against a discrete DQN baseline. DQN completely failed because with 7 joints, the action combinations exceed 2,000—it's a search-space nightmare. Dense rewards also failed by inducing 'safe hovering'. Our Sparse PPO is the clear mechanical winner.",
    [
        {'text': 'INSERT: sparse_dqn_6k.png', 'left': Inches(5.0), 'top': Inches(1.5), 'width': Inches(4.5), 'height': Inches(2.5)},
        {'text': 'INSERT: dense_ppo_10k.png', 'left': Inches(5.0), 'top': Inches(4.5), 'width': Inches(4.5), 'height': Inches(2.5)}
    ]
)

# --- SLIDE 10: Evidence ---
add_content_slide(
    prs,
    "Cross-Morphology Performance Matrix",
    [
        "7-DoF model outperformed the 3-DoF model in convergence reliability.",
        "Kinematic redundancy ($n > 6$) provides a smoother gradient for implicit IK learning.",
        "Demonstrates robustness against varying inertial loads and friction."
    ],
    "This table is the validation of the 'Adaptive' claim. Notice that the 7-DOF Panda actually performed better than the 3-DOF. This suggests that Kinematic Redundancy helps the RL agent find more viable paths to the goal, making it easier to solve than a restricted 3-DOF system.",
    [{'text': 'INSERT: Table\nDoF | Success Rate | Time\n3-DoF | 42% | 2.8h\n5-DoF | 38% | 3.2h\n7-DoF | 51% | 4.5h', 'left': Inches(5.2), 'top': Inches(2.0), 'width': Inches(4.5), 'height': Inches(3.5)}]
)

# --- SLIDE 11: Impact ---
add_content_slide(
    prs,
    "Future Work: Morphology Co-Optimization",
    [
        "Domain randomization for sim-to-real transfer on physical hardware.",
        "Parameterizing link lengths and joint offsets within the optimization loop.",
        "Enabling robots to evolve physical structures alongside control policies."
    ],
    "The future is Co-Optimization. Now that the brain can adapt to any body, we can use the brain to improve the body. We can let the AI suggest optimal link lengths and joint placements to maximize success rates, leading to robots that design themselves for specific tasks.",
    [{'text': 'INSERT: Robot Evolution / Co-Optimization Diagram', 'left': Inches(5.2), 'top': Inches(2.0), 'width': Inches(4.5), 'height': Inches(3.5)}]
)

# --- SLIDE 12: Impact ---
add_content_slide(
    prs,
    "Conclusion and Synthesis",
    [
        "Developed a structure-agnostic pipeline for automated IK synthesis.",
        "Proved the efficacy of massive sparse rewards in contact-rich physics.",
        "Successfully achieved a 51% implicit IK success rate on a 7-DoF manipulator."
    ],
    "In conclusion, we have established a robust, structure-agnostic pipeline for robotic control synthesis. By proving the efficacy of massive sparse rewards across varying DoF, we have lowered the technical barrier for custom robot development. I am now open to technical questions.",
    [{'text': 'INSERT: Panda_ppo_sparse_10k.png\n(Terminal 51% Screenshot)', 'left': Inches(5.2), 'top': Inches(2.0), 'width': Inches(4.5), 'height': Inches(3.5)}]
)

prs.save('Defense_Presentation.pptx')
print("Presentation generated successfully: Defense_Presentation.pptx")